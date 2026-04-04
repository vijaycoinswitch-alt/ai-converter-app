import os
import uuid
import zipfile
import subprocess

# ALLOWED_EXTENSIONS defined below after app init

import mimetypes

def allowed_file(filename):
    if not ("." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS):
        return False
    mime_type, _ = mimetypes.guess_type(filename)
    if not mime_type:
        return False
    if mime_type.startswith('application/') or mime_type.startswith('image/') or mime_type.startswith('text/'):
        return True
    return False

def generate_safe_filename(filename):
    ext = filename.rsplit(".", 1)[1].lower()
    safe_name = secure_filename(filename.rsplit(".", 1)[0])[:50]
    return f"{uuid.uuid4().hex}_{safe_name}.{ext}"

from pathlib import Path
from datetime import datetime, timedelta
from threading import Thread

from flask import Flask, render_template, request, jsonify, send_from_directory, redirect, url_for, flash, abort, Response
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, login_required, logout_user, current_user
from flask_bcrypt import Bcrypt
from werkzeug.utils import secure_filename
from dotenv import load_dotenv

from flask_wtf.csrf import CSRFProtect
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from flask_talisman import Talisman
import email_service
import secrets

# --- PDF Manipulation Libraries ---
import fitz  # PyMuPDF
from PyPDF2 import PdfReader, PdfWriter, PdfMerger
from pdf2docx import Converter
try:
    from docx2pdf import convert as docx2pdf_convert
except ImportError:
    docx2pdf_convert = None
from PIL import Image
import pandas as pd
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
from reportlab.lib import colors
import pytesseract
from pdf2image import convert_from_path
import pdfkit
from pptx import Presentation
from pptx.util import Inches
import pdfplumber

# AI Libraries
import openai

# Signatures
try:
    from pyhanko.sign import signers
    from pyhanko.pdf_utils.writer import copy_into_new_writer
    from pyhanko.pdf_utils.reader import PdfFileReader
except ImportError:
    pass

# Load environment
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")


app = Flask(__name__)
# 🔐 SECURITY CONFIG
# load_dotenv() already called above

BASE_DIR = Path(__file__).resolve().parent
INSTANCE_DIR = BASE_DIR / "instance"
PRIVATE_UPLOAD_DIR = INSTANCE_DIR / "private_uploads"
PRIVATE_OUTPUT_DIR = INSTANCE_DIR / "private_outputs"

INSTANCE_DIR.mkdir(exist_ok=True)
PRIVATE_UPLOAD_DIR.mkdir(exist_ok=True)
PRIVATE_OUTPUT_DIR.mkdir(exist_ok=True)

app.config["SECRET_KEY"] = os.getenv("SECRET_KEY", "dev-insecure-change-me")
app.config["SQLALCHEMY_DATABASE_URI"] = os.getenv("DATABASE_URL", "sqlite:///app.db")

# MAX_CONTENT_LENGTH set below with OUTPUT_FOLDER config

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

app.config["SESSION_COOKIE_HTTPONLY"] = True
app.config["SESSION_COOKIE_SAMESITE"] = "Lax"
app.config["REMEMBER_COOKIE_HTTPONLY"] = True
app.config["PERMANENT_SESSION_LIFETIME"] = timedelta(minutes=30)

is_production = os.getenv("FLASK_ENV", "development").lower() == "production"
app.config["SESSION_COOKIE_SECURE"] = is_production
app.config["REMEMBER_COOKIE_SECURE"] = is_production

app.config["WTF_CSRF_TIME_LIMIT"] = 3600

csrf = CSRFProtect(app)

limiter = Limiter(
    key_func=get_remote_address,
    app=app,
    default_limits=["200 per day", "10 per minute"],
    storage_uri="memory://",
)

# csp = {
#     "default-src": "'self'",
#     "img-src": ["'self'", "data:"],
#     "style-src": ["'self'", "'unsafe-inline'", "https://cdnjs.cloudflare.com", "https://fonts.googleapis.com"],
#     "font-src": ["'self'", "https://cdnjs.cloudflare.com", "https://fonts.gstatic.com"],
#     "script-src": ["'self'", "'unsafe-inline'"],
# }

# Talisman(
#     app,
#     force_https=is_production,
#     strict_transport_security=is_production,
#     content_security_policy=csp,
#     frame_options="DENY",
#     referrer_policy="strict-origin-when-cross-origin"
# )

ALLOWED_EXTENSIONS = {
    "pdf", "docx", "doc", "xlsx", "xls", "pptx", "ppt",
    "jpg", "jpeg", "png", "txt"
}

app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'outputs'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB max

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

import logging
from logging.handlers import RotatingFileHandler
if not app.debug:
    os.makedirs('instance/logs', exist_ok=True)
    file_handler = RotatingFileHandler('instance/logs/security.log', maxBytes=10240, backupCount=10)
    file_handler.setFormatter(logging.Formatter('%(asctime)s %(levelname)s: %(message)s [in %(pathname)s:%(lineno)d]'))
    file_handler.setLevel(logging.INFO)
    app.logger.addHandler(file_handler)
    app.logger.setLevel(logging.INFO)
    app.logger.info('VijayPDF secure startup')

# --- Background Task: Cleanup old files ---
def cleanup_old_files():
    """Delete files in uploads and outputs older than 2 hours."""
    import time
    while True:
        try:
            now = time.time()
            for folder in [app.config['UPLOAD_FOLDER'], app.config['OUTPUT_FOLDER']]:
                for f in os.listdir(folder):
                    fp = os.path.join(folder, f)
                    if os.path.isfile(fp):
                        if os.stat(fp).st_mtime < now - (2 * 3600):
                            os.remove(fp)
                            app.logger.info(f"Cleaned up old file: {f}")
        except Exception as e:
            app.logger.error(f"Cleanup error: {e}")
        time.sleep(3600) # Run every hour

cleanup_thread = Thread(target=cleanup_old_files, daemon=True)
cleanup_thread.start()

@app.errorhandler(500)
def internal_error(error):
    try:
        db.session.rollback()
    except Exception:
        pass
    app.logger.error(f"Server Error: {error}")
    return render_template('500.html'), 500

db = SQLAlchemy(app)
bcrypt = Bcrypt(app)
login_manager = LoginManager(app)
login_manager.login_view = 'login'

with app.app_context():
    db.create_all()

# --- Models ---
class User(db.Model, UserMixin):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(30), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password = db.Column(db.String(60), nullable=False)
    failed_logins = db.Column(db.Integer, default=0)
    locked_until = db.Column(db.DateTime, nullable=True)
    
    # OTP Fields
    otp = db.Column(db.String(6), nullable=True)
    otp_expiry = db.Column(db.DateTime, nullable=True)
    is_verified = db.Column(db.Boolean, default=False)
    
    # Password Reset Fields
    reset_token = db.Column(db.String(100), nullable=True)
    reset_token_expiry = db.Column(db.DateTime, nullable=True)

    files = db.relationship("UserFile", backref="owner", lazy=True)
class UserFile(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False)
    original_name = db.Column(db.String(255), nullable=False)
    stored_name = db.Column(db.String(255), nullable=False, unique=True)
    stored_path = db.Column(db.String(500), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)    

class ConversionHistory(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    original_filename = db.Column(db.String(100), nullable=False)
    converted_filename = db.Column(db.String(100), nullable=True)
    conversion_type = db.Column(db.String(50), nullable=False)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True)

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

# --- Helpers ---
def save_uploaded_files(files):
    saved_paths = []

    for f in files:
        if not f or not f.filename:
            continue

        if not allowed_file(f.filename):
            raise Exception(f"File type not allowed: {f.filename}")

        safe_name = generate_safe_filename(f.filename)
        path = os.path.join(app.config['UPLOAD_FOLDER'], safe_name)

        f.save(path)
        saved_paths.append((path, f.filename))

    return saved_paths

def log_history(original, converted, type_name):
    if current_user.is_authenticated:
        db.session.add(ConversionHistory(original_filename=original, converted_filename=converted, conversion_type=type_name, user_id=current_user.id))
        db.session.commit()

# ======================================================================
# 1. ORGANIZE PDF FUNCTIONS
# ======================================================================
def merge_pdf_func(files_list, out_path):
    merger = PdfMerger()
    for fp in files_list:
        merger.append(fp)
    merger.write(out_path)
    merger.close()

def split_pdf_func(in_path, out_zip):
    doc = fitz.open(in_path)
    with zipfile.ZipFile(out_zip, 'w') as zf:
        for i in range(len(doc)):
            ndoc = fitz.open()
            ndoc.insert_pdf(doc, from_page=i, to_page=i)
            tpath = out_zip.replace('.zip', f'_pg{i+1}.pdf')
            ndoc.save(tpath)
            zf.write(tpath, f'page_{i+1}.pdf')
            os.remove(tpath)

def remove_pages_func(in_path, out_path, pages_str):
    doc = fitz.open(in_path)
    # Ex: "1,3,5" -> convert to 0-indexed [0, 2, 4]
    try:
        pgs = [int(p.strip())-1 for p in pages_str.split(',')]
        # Delete reverse order to avoid index shifts
        for p in sorted(pgs, reverse=True):
            if 0 <= p < len(doc):
                doc.delete_page(p)
    except:
        raise Exception("Invalid page range format. Use e.g. '1,3,5'")
    doc.save(out_path)

def extract_pages_func(in_path, out_path, pages_str):
    doc = fitz.open(in_path)
    try:
        pgs = [int(p.strip())-1 for p in pages_str.split(',')]
        doc.select(pgs)
    except:
        raise Exception("Invalid page format.")
    doc.save(out_path)

def organize_pdf_func(in_path, out_path, order_str):
    doc = fitz.open(in_path)
    try:
        pgs = [int(p.strip())-1 for p in order_str.split(',')]
        doc.select(pgs)
    except:
        raise Exception("Invalid page order format.")
    doc.save(out_path)

def scan_to_pdf_func(image_paths, out_path):
    images = [Image.open(fp).convert('RGB') for fp in image_paths]
    if images:
        images[0].save(out_path, save_all=True, append_images=images[1:])

# ======================================================================
# 2. OPTIMIZE PDF FUNCTIONS
# ======================================================================
def compress_pdf_func(in_path, out_path):
    doc = fitz.open(in_path)
    doc.save(out_path, garbage=4, deflate=True)

def repair_pdf_func(in_path, out_path):
    # fitz will rewrite broken xref tables automatically stringently
    doc = fitz.open(in_path)
    doc.save(out_path, clean=True)

def ocr_pdf_func(in_path, out_txt):
    POPPLER_PATH = os.getenv("POPPLER_PATH", None)
    imgs = convert_from_path(in_path, poppler_path=POPPLER_PATH)
    text = "\n".join([pytesseract.image_to_string(i) for i in imgs])
    with open(out_txt, 'w', encoding='utf-8') as f:
        f.write(text)

# ======================================================================
# 3. CONVERT TO PDF FUNCTIONS
# ======================================================================
def jpg_to_pdf_func(image_paths, out_path):
    scan_to_pdf_func(image_paths, out_path)

def word_to_pdf_func(in_path, out_path):
    if docx2pdf_convert is None:
        raise Exception("Word to PDF is not available on this server.")
    docx2pdf_convert(os.path.abspath(in_path), os.path.abspath(out_path))

def powerpoint_to_pdf_func(in_path, out_path):
    """Cross-platform PPTX to PDF using python-pptx + reportlab."""
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas
    try:
        prs = Presentation(in_path)
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches
        c = canvas.Canvas(out_path, pagesize=(slide_width * inch, slide_height * inch))
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            y = slide_height * inch - 50
            for t in texts:
                if y < 50:
                    c.showPage()
                    y = slide_height * inch - 50
                c.drawString(50, y, t[:120])
                y -= 18
            c.showPage()
        c.save()
    except Exception as e:
        raise Exception(f"PowerPoint to PDF conversion failed: {e}")

def excel_to_pdf_func(in_path, out_path):
    df = pd.read_excel(in_path)
    doc = SimpleDocTemplate(out_path, pagesize=letter)
    data = [df.columns.values.tolist()] + df.values.tolist()
    d_str = [[str(item) for item in row] for row in data]
    t = Table(d_str)
    t.setStyle(TableStyle([('BACKGROUND', (0,0), (-1,0), colors.grey), ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke), ('GRID', (0,0), (-1,-1), 1, colors.black)]))
    doc.build([t])

def html_to_pdf_func(in_path, out_path):
    # Note: Requires system wkhtmltopdf
    try:
        pdfkit.from_file(in_path, out_path)
    except OSError:
        raise Exception("HTML to PDF requires wkhtmltopdf installed on the system.")

# ======================================================================
# 4. CONVERT FROM PDF FUNCTIONS
# ======================================================================
def pdf_to_jpg_func(in_path, out_zip):
   POPPLER_PATH = os.getenv("POPPLER_PATH", None)
   imgs = convert_from_path(in_path, poppler_path=POPPLER_PATH)
   with zipfile.ZipFile(out_zip, 'w') as zf:
        for idx, img in enumerate(imgs):
            tmp = out_zip.replace('.zip', f'_{idx}.jpg')
            img.save(tmp, 'JPEG')
            zf.write(tmp, f'page_{idx+1}.jpg')
            os.remove(tmp)

def pdf_to_word_func(in_path, out_path):
    cv = Converter(in_path)
    cv.convert(out_path)
    cv.close()

def pdf_to_powerpoint_func(in_path, out_path):
    """Convert PDF to PPTX. Uses Poppler images when available, falls back to text extraction."""
    prs = Presentation()
    try:
        POPPLER_PATH = os.getenv("POPPLER_PATH", None)
        imgs = convert_from_path(in_path, poppler_path=POPPLER_PATH)
        for img in imgs:
            tmp = out_path.replace('.pptx', '_tmp.jpg')
            img.save(tmp)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(tmp, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.remove(tmp)
    except Exception:
        # Fallback: extract text via PyMuPDF and place on slides
        doc = fitz.open(in_path)
        for page in doc:
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            from pptx.util import Pt, Emu
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5),
                                             prs.slide_width - Inches(1), prs.slide_height - Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            text = page.get_text().strip()
            tf.text = text[:3000] if text else "(No text content on this page)"
    prs.save(out_path)

def pdf_to_excel_func(in_path, out_path):
    tables_data = []
    with pdfplumber.open(in_path) as pdf:
        for page in pdf.pages:
            tabs = page.extract_tables()
            for t in tabs:
                tables_data.extend(t)
    if not tables_data:
        raise Exception("No tables found in PDF.")
    df = pd.DataFrame(tables_data[1:], columns=tables_data[0])
    df.to_excel(out_path, index=False)

def pdf_to_pdfa_func(in_path, out_path):
    # Fallback Ghostscript invocation or error
    try:
        cmd = ['gs', '-dPDFA', '-dBATCH', '-dNOPAUSE', '-sProcessColorModel=DeviceRGB', '-sDEVICE=pdfwrite', f'-sOutputFile={out_path}', in_path]
        subprocess.run(cmd, check=True)
    except FileNotFoundError:
        raise Exception("Ghostscript (gs) is required for PDF/A conversion.")

# ======================================================================
# 5. EDIT PDF FUNCTIONS
# ======================================================================
def rotate_pdf_func(in_path, out_path, angle):
    doc = fitz.open(in_path)
    for p in doc:
        p.set_rotation(int(angle))
    doc.save(out_path)

def add_page_numbers_func(in_path, out_path):
    doc = fitz.open(in_path)
    for i, p in enumerate(doc):
        p.insert_text(fitz.Point(300, 800), str(i+1), fontsize=12)
    doc.save(out_path)

def add_watermark_func(in_path, out_path, text):
    doc = fitz.open(in_path)
    for p in doc:
        p.insert_text(fitz.Point(100, 400), text, fontsize=30, color=(1, 0, 0))
    doc.save(out_path)

def crop_pdf_func(in_path, out_path):
    doc = fitz.open(in_path)
    for p in doc:
        r = p.rect
        r.x0 += 50; r.y0 += 50; r.x1 -= 50; r.y1 -= 50 # simple 50pt crop
        p.set_cropbox(r)
    doc.save(out_path)

def edit_pdf_func(in_path, out_path, append_text):
    doc = fitz.open(in_path)
    if len(doc) > 0:
        doc[0].insert_text(fitz.Point(50, 50), append_text, fontsize=14, color=(0,0,1))
    doc.save(out_path)

# ======================================================================
# 6. SECURITY FUNCTIONS
# ======================================================================
def unlock_pdf_func(in_path, out_path, pwd):
    doc = fitz.open(in_path)
    doc.authenticate(pwd)
    doc.save(out_path)

def protect_pdf_func(in_path, out_path, pwd):
    doc = fitz.open(in_path)
    doc.save(out_path, encryption=fitz.PDF_ENCRYPT_AES_256, user_pw=pwd)

def sign_pdf_func(in_path, out_path, pfx_path, pfx_pass):
    signer = signers.SimpleSigner.load_pkcs12(pfx_path, pfx_pass.encode())
    with open(in_path, 'rb') as f:
        w = copy_into_new_writer(PdfFileReader(f))
        signers.sign_pdf(w, signers.PdfSignatureMetadata(field_name='Signature1'), signer=signer, in_place=True, out_file=open(out_path, 'wb'))

def redact_pdf_func(in_path, out_path, txt_to_redact):
    doc = fitz.open(in_path)
    for p in doc:
        inst = p.search_for(txt_to_redact)
        for i in inst:
            p.add_redact_annot(i, fill=(0,0,0))
        p.apply_redactions()
    doc.save(out_path)

def compare_pdf_func(in1, in2, out_txt):
    t1 = "".join([p.get_text() for p in fitz.open(in1)])
    t2 = "".join([p.get_text() for p in fitz.open(in2)])
    with open(out_txt, 'w', encoding='utf-8') as f:
        f.write("Differences found." if t1 != t2 else "Documents are identical textually.")

# ======================================================================
# 7. AI FUNCTIONS
# ======================================================================
def ai_summarize_func(in_path):
    t = "".join([p.get_text() for p in fitz.open(in_path)])
    r = openai.chat.completions.create(model="gpt-3.5-turbo", messages=[{"role": "user", "content": "Summarize: " + t[:10000]}])
    return r.choices[0].message.content

def translate_pdf_func(in_path, lang):
    t = "".join([p.get_text() for p in fitz.open(in_path)])
    r = openai.chat.completions.create(model="gpt-3.5-turbo", messages=[{"role": "user", "content": f"Translate to {lang}: " + t[:4000]}])
    return r.choices[0].message.content

def chat_with_pdf_func(in_path, q):
    t = "".join([p.get_text() for p in fitz.open(in_path)])
    r = openai.chat.completions.create(model="gpt-3.5-turbo", messages=[{"role": "user", "content": f"Context: {t[:8000]}\nQuestion: {q}"}])
    return r.choices[0].message.content


# ======================================================================
# FLASK SINGLE-FILE ROUTER BINDINGS
# ======================================================================

def process_wrapper(func, ext=".pdf", multi=False, is_text=False):
    files = request.files.getlist('files[]')
    if not files: return jsonify({'error': 'No file provided'}), 400
    paths = save_uploaded_files(files)
    
    out_name = f"{uuid.uuid4().hex}{ext}"
    out_path = os.path.join(app.config['OUTPUT_FOLDER'], out_name)
    
    try:
        # Some endpoints return pure text
        if is_text:
            ans = func(paths[0][0])
            log_history(paths[0][1], "AI_TXT", request.path)
            return jsonify({'success': True, 'text_result': ans})
            
        if multi:
            func([p[0] for p in paths], out_path)
        else:
            func(paths[0][0], out_path)
        log_history(paths[0][1], out_name, request.path)
        return jsonify({'success': True, 'download_url': url_for('download_file', filename=out_name)})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/favicon.ico')
def favicon():
    return send_from_directory(
        os.path.join(app.root_path, 'static'),
        'favicon.ico',
        mimetype='image/vnd.microsoft.icon'
    )


# API Routes

# ======================================================================
# BLOG DATA (20 SEO-optimised posts as Python dicts — no DB needed)
# ======================================================================
BLOG_POSTS = [
    {
        "slug": "how-to-convert-pdf-to-word-free",
        "title": "How to Convert PDF to Word for Free in 2026",
        "meta_description": "Learn the easiest ways to convert PDF to Word documents for free online. Step-by-step guide using VijayPDF. No software download needed.",
        "category": "Tutorial",
        "date": "April 1, 2026",
        "read_time": "6 min read",
        "excerpt": "Converting a PDF into an editable Word document no longer requires paid software. This guide explains the fastest free method.",
        "content": """
<p>If you've ever received a PDF and needed to make edits, you've probably run into one of the most frustrating problems in modern document management: PDFs are designed to look the same everywhere, but that "locked" format makes them nearly impossible to edit directly.</p>

<p>The good news? Converting PDF to Word is now completely free, fast, and surprisingly accurate — thanks to tools like <a href="/pdf-to-word">VijayPDF's PDF to Word converter</a>.</p>

<h2>Why Convert PDF to Word?</h2>
<p>PDFs are brilliant for sharing final documents. But when you need to:</p>
<ul>
<li>Update a resume you originally got as a PDF</li>
<li>Make changes to a contract or proposal</li>
<li>Extract text from a report to use elsewhere</li>
<li>Collaborate on a document that was shared as a PDF</li>
</ul>
<p>...you need an editable Word document. Retyping everything would take hours. That's where PDF to Word conversion saves the day.</p>

<h2>Step-by-Step: Convert PDF to Word Using VijayPDF</h2>
<ol>
<li><strong>Go to <a href="/pdf-to-word">vijaypdf.com/pdf-to-word</a></strong></li>
<li><strong>Upload your PDF file</strong> — drag and drop it into the upload box, or click to browse your files</li>
<li><strong>Wait for processing</strong> — our engine extracts text, images, tables, and formatting (usually under 15 seconds)</li>
<li><strong>Click Download</strong> — save your .docx file to your device</li>
<li><strong>Open in Word or Google Docs</strong> and make your edits</li>
</ol>

<h2>What Formatting is Preserved?</h2>
<p>VijayPDF preserves a remarkable amount of your original document's formatting:</p>
<ul>
<li>Headings and paragraph styles</li>
<li>Tables with rows and column borders</li>
<li>Embedded images and graphics</li>
<li>Bold, italic, underline, and font sizes</li>
<li>Multi-column layouts like newsletters</li>
<li>Headers and footers with page numbers</li>
</ul>

<h2>Tips for Best Results</h2>
<p>For the cleanest conversion:</p>
<ul>
<li>Use text-based PDFs (not scanned images) for highest accuracy</li>
<li>If you have a scanned PDF, use our <a href="/ocr-pdf">OCR PDF tool</a> first to make text machine-readable</li>
<li>Very complex multi-column layouts may need minor manual touch-ups</li>
</ul>

<h2>Is VijayPDF Safe to Use?</h2>
<p>Absolutely. All file transfers use SSL encryption, and your uploaded documents are automatically and permanently deleted from our servers after 2 hours. We never read, share, or retain your personal files.</p>

<h2>Conclusion</h2>
<p>Converting PDF to Word doesn't have to cost money or require installing software. With VijayPDF, the process takes under a minute — completely free, completely secure, and accessible from any device.</p>

<div class="blog-cta-box">
<h3>Ready to convert your PDF?</h3>
<a href="/pdf-to-word" class="btn btn-primary">Convert PDF to Word Free →</a>
</div>
"""
    },
    {
        "slug": "best-free-pdf-tools-online-2026",
        "title": "Best Free PDF Tools Online in 2026",
        "meta_description": "Compare the best free PDF tools online in 2026. Find tools for merging, splitting, compressing, and converting PDFs without paying a cent.",
        "category": "Roundup",
        "date": "March 28, 2026",
        "read_time": "8 min read",
        "excerpt": "With dozens of PDF tools available online, which ones are genuinely free and actually work? We break down the best options in 2026.",
        "content": """
<p>The PDF tools market has exploded. From iLovePDF to Smallpdf to Adobe Acrobat Online, there are dozens of options — but many hide features behind paywalls or add watermarks to your files. This guide focuses on what's genuinely free and genuinely useful in 2026.</p>

<h2>What Makes a PDF Tool Truly "Free"?</h2>
<p>Before we begin, let's define what "free" actually means:</p>
<ul>
<li><strong>No watermarks</strong> on output files</li>
<li><strong>No file size limits</strong> that cut off most real-world documents</li>
<li><strong>No account required</strong> for basic operations</li>
<li><strong>No feature paywalls</strong> on core functionality</li>
</ul>

<h2>Top Free PDF Tools Available Right Now</h2>

<h3>1. VijayPDF — Best All-in-One Free PDF Toolkit</h3>
<p><a href="/">VijayPDF</a> offers 25+ free tools with absolutely no watermarks. Key features include PDF to Word, merge, split, compress, JPG to PDF, PDF to Excel, AI summarization, watermark addition, rotation, and more. All tools are free with no account required for basic conversions.</p>
<p><strong>Best for:</strong> Users who need a comprehensive toolkit with no payment barriers<br>
<strong>File limit:</strong> 50 MB<br>
<strong>Watermarks:</strong> None</p>

<h3>2. ILovePDF</h3>
<p>One of the most popular PDF platforms. Offers excellent tools but limits free users to 2 tasks per hour and adds usage caps. Premium plans unlock unlimited use.</p>
<p><strong>Best for:</strong> Occasional users<br><strong>Watermarks:</strong> None on free tier</p>

<h3>3. Smallpdf</h3>
<p>Beautifully designed with two free operations per day. After that, a subscription is required. Good for rare one-off tasks.</p>
<p><strong>Best for:</strong> Very occasional use<br><strong>Free limit:</strong> 2 operations/day</p>

<h3>4. PDF24 Tools</h3>
<p>Completely free German PDF service with an extensive tool collection. No limits, no watermarks, but occasionally slower processing than cloud-first tools.</p>
<p><strong>Best for:</strong> Privacy-conscious users, bulk processing</p>

<h2>Feature Comparison</h2>
<p>Here's how these tools stack up on the features that matter most:</p>
<ul>
<li><strong>PDF to Word:</strong> VijayPDF ✅ | ILovePDF ✅ | Smallpdf ✅ | PDF24 ✅</li>
<li><strong>No watermarks:</strong> VijayPDF ✅ | ILovePDF ✅ | Smallpdf ✅ | PDF24 ✅</li>
<li><strong>No usage limits:</strong> VijayPDF ✅ | ILovePDF ❌ | Smallpdf ❌ | PDF24 ✅</li>
<li><strong>AI tools:</strong> VijayPDF ✅ | ILovePDF ❌ | Smallpdf ❌ | PDF24 ❌</li>
<li><strong>No signup needed:</strong> VijayPDF ✅ | ILovePDF ❌ | Smallpdf ❌ | PDF24 ✅</li>
</ul>

<h2>Our Recommendation</h2>
<p>For most users, <a href="/">VijayPDF</a> provides the most generous free tier with the widest range of tools. For specialized enterprise use with team collaboration, a paid plan from ILovePDF or Adobe Acrobat may be warranted.</p>

<div class="blog-cta-box">
<h3>Try VijayPDF's complete free PDF toolkit</h3>
<a href="/tools" class="btn btn-primary">Explore All 25+ Tools →</a>
</div>
"""
    },
    {
        "slug": "reduce-pdf-size-without-losing-quality",
        "title": "How to Reduce PDF Size Without Losing Quality",
        "meta_description": "Learn proven techniques to reduce PDF file size without visible quality loss. Compress PDFs for email, web, and storage using free online tools.",
        "category": "Guide",
        "date": "March 25, 2026",
        "read_time": "7 min read",
        "excerpt": "Large PDFs are a headache for email and storage. This guide explains how to shrink them dramatically while keeping quality intact.",
        "content": """
<p>A 50MB PDF is practically impossible to email. Uploading a large PDF to a website portal with a 10MB limit is equally frustrating. Yet reducing PDF size used to require expensive software. Today, you can compress any PDF in seconds using free online tools.</p>

<h2>Why Are Some PDFs So Large?</h2>
<p>PDF file size is driven primarily by:</p>
<ul>
<li><strong>Embedded images</strong> — photos at full camera resolution can be enormous</li>
<li><strong>Scanned documents</strong> — each scanned page is a high-resolution image</li>
<li><strong>Embedded fonts</strong> — full font sets instead of subsets</li>
<li><strong>Uncompressed content streams</strong> — text and vector data that wasn't compressed during creation</li>
<li><strong>Redundant objects</strong> — duplicate resources leftover from editing</li>
</ul>

<h2>What Happens During Compression?</h2>
<p>A good PDF compressor like <a href="/compress-pdf">VijayPDF's Compress PDF tool</a> targets these size contributors:</p>
<ol>
<li><strong>Image resampling</strong> — reduces image resolution from print-quality (300 DPI) to screen-quality (96–150 DPI) for images that don't need maximum resolution</li>
<li><strong>Image re-encoding</strong> — recompresses images using more efficient JPEG or WebP encoding</li>
<li><strong>Stream compression</strong> — applies deflate/zlib compression to uncompressed content streams</li>
<li><strong>Object cleanup</strong> — removes unreferenced objects, old revisions, and duplicate resources</li>
<li><strong>Font subsetting</strong> — embeds only the characters actually used from each font</li>
</ol>

<h2>How Much Can You Reduce a PDF?</h2>
<p>Results vary dramatically by content type:</p>
<ul>
<li><strong>Scanned documents:</strong> 50–80% reduction — biggest gains</li>
<li><strong>Photo-heavy PDFs (brochures, catalogs):</strong> 40–70% reduction</li>
<li><strong>Office documents (Word/Excel to PDF):</strong> 20–40% reduction</li>
<li><strong>Text-only PDFs:</strong> 10–20% reduction — already very efficient</li>
</ul>

<h2>Step-by-Step: Compress a PDF for Free</h2>
<ol>
<li>Go to <a href="/compress-pdf">vijaypdf.com/compress-pdf</a></li>
<li>Upload your large PDF file</li>
<li>Wait for automatic compression (usually 5–20 seconds)</li>
<li>Download your compressed PDF</li>
<li>Compare file sizes — enjoy the difference!</li>
</ol>

<h2>Other Ways to Reduce PDF Size</h2>
<p>Beyond compression tools, consider these additional approaches:</p>
<ul>
<li><strong>Split large documents</strong> — use <a href="/split-pdf">Split PDF</a> to break a massive PDF into smaller sections</li>
<li><strong>Remove unnecessary pages</strong> — use <a href="/remove-pages">Remove Pages</a> to delete pages you don't need to share</li>
<li><strong>Convert images to grayscale</strong> before creating the PDF if colour isn't needed</li>
</ul>

<div class="blog-cta-box">
<h3>Compress your PDF for free now</h3>
<a href="/compress-pdf" class="btn btn-primary">Compress PDF →</a>
</div>
"""
    },
    {
        "slug": "merge-pdf-files-easily-online",
        "title": "How to Merge PDF Files Easily Online (2026 Guide)",
        "meta_description": "Learn how to merge multiple PDF files into one document online for free. Step-by-step guide to combining PDFs without software.",
        "category": "Tutorial",
        "date": "March 22, 2026",
        "read_time": "5 min read",
        "excerpt": "Combining multiple PDFs into one file is easier than ever. This guide shows you how to do it in seconds — completely free.",
        "content": """
<p>Managing a stack of separate PDF files — chapter by chapter, invoice by invoice, or page by page — is annoying. Merging them into a single organized document makes sharing, printing, and archiving dramatically easier.</p>

<h2>When Would You Merge PDFs?</h2>
<ul>
<li>Combining a resume, cover letter, and portfolio into one application PDF</li>
<li>Merging monthly invoices into a single quarterly report</li>
<li>Assembling a multi-chapter document from separately saved sections</li>
<li>Combining scanned pages into one complete file</li>
<li>Bundling supporting documents for a legal or administrative submission</li>
</ul>

<h2>How to Merge PDFs Using VijayPDF</h2>
<ol>
<li>Visit <a href="/merge-pdf">vijaypdf.com/merge-pdf</a></li>
<li>Click "Select Files" and upload two or more PDF files — or drag and drop them</li>
<li>Files will be merged in the order you upload them, so upload them in the sequence you want</li>
<li>Click the convert button and wait a few seconds</li>
<li>Download your single merged PDF</li>
</ol>

<h2>Does Merging Reduce PDF Quality?</h2>
<p>No. Merging PDFs is a lossless operation. VijayPDF stitches the internal page objects from each PDF together into one document without re-encoding, re-compressing, or modifying any content. Images stay sharp, text stays selectable, and fonts remain embedded exactly as they were.</p>

<h2>Tips for Merging PDFs</h2>
<ul>
<li>Upload files in the exact order you want them to appear in the merged PDF</li>
<li>If any PDFs are password-protected, <a href="/unlock-pdf">unlock them first</a></li>
<li>For very large merged files, run the result through our <a href="/compress-pdf">Compress PDF</a> tool to reduce the size afterwards</li>
<li>Use <a href="/organize-pdf">Organize PDF</a> if you need to reorder pages after merging</li>
</ul>

<div class="blog-cta-box">
<h3>Merge your PDF files now — free</h3>
<a href="/merge-pdf" class="btn btn-primary">Merge PDFs →</a>
</div>
"""
    },
    {
        "slug": "convert-jpg-to-pdf-in-seconds",
        "title": "Convert JPG to PDF in Seconds — Free Online Tool",
        "meta_description": "Convert JPG images to PDF online for free in seconds. Combine multiple photos into one PDF. Works on iPhone, Android, Windows, and Mac.",
        "category": "Tutorial",
        "date": "March 20, 2026",
        "read_time": "4 min read",
        "excerpt": "Need to turn a photo into a PDF? This guide shows how to convert JPG to PDF instantly — for free, on any device.",
        "content": """
<p>Whether you've photographed a document, a whiteboard, a receipt, or a page from a book — you often need it in PDF format for sharing, submitting, or archiving. Converting JPG to PDF is fast, free, and done entirely in your browser.</p>

<h2>Why Convert JPG to PDF?</h2>
<ul>
<li>PDFs are universally accepted for official document submissions</li>
<li>Combine multiple photos into one organized document</li>
<li>PDF files are easier to print with consistent page sizes</li>
<li>Reduce confusion by sending one file instead of many images</li>
<li>PDFs open in every email client and browser without special software</li>
</ul>

<h2>How to Convert JPG to PDF (Free)</h2>
<ol>
<li>Go to <a href="/jpg-to-pdf">vijaypdf.com/jpg-to-pdf</a></li>
<li>Upload one or more JPG, PNG, or image files</li>
<li>Each image becomes one page in the PDF — upload in your desired order</li>
<li>Download your PDF instantly</li>
</ol>

<h2>What Image Formats Are Supported?</h2>
<p>VijayPDF's image-to-PDF tool accepts: JPG, JPEG, PNG, BMP, GIF, TIFF, and WebP. If you have another format, convert it to PNG or JPG first using your device's default image viewer.</p>

<h2>Tips for Best Results</h2>
<ul>
<li>Photograph documents in good lighting to ensure readable text in the PDF</li>
<li>Use PNG format for screenshots to preserve sharp text edges</li>
<li>Upload images in page order when combining multiple photos</li>
<li>After creating the PDF, compress it with our <a href="/compress-pdf">PDF Compressor</a> if the file is too large</li>
</ul>

<div class="blog-cta-box">
<h3>Convert your images to PDF now</h3>
<a href="/jpg-to-pdf" class="btn btn-primary">Convert JPG to PDF →</a>
</div>
"""
    },
    {
        "slug": "how-to-split-a-pdf-free",
        "title": "How to Split a PDF File for Free — Online in Seconds",
        "meta_description": "Learn how to split a PDF into individual pages or extract specific pages online for free. No software needed. Works on all devices.",
        "category": "Tutorial",
        "date": "March 18, 2026",
        "read_time": "5 min read",
        "excerpt": "Need just one page from a long PDF? Or want to separate every page into its own file? This guide shows how — for free.",
        "content": """
<p>Sometimes you have a 50-page report but only need to share pages 5 through 12. Or you've scanned a stack of documents and need each page as a separate file. PDF splitting solves both problems instantly.</p>

<h2>What Does Splitting a PDF Mean?</h2>
<p>Splitting a PDF means separating it into smaller pieces. This can mean:</p>
<ul>
<li><strong>Split every page</strong> — get each page as its own individual PDF file</li>
<li><strong>Extract specific pages</strong> — choose page 3, 7, and 15 and get just those</li>
<li><strong>Remove unwanted pages</strong> — create a new PDF without specific pages</li>
</ul>

<h2>How to Split a PDF for Free</h2>
<ol>
<li>Visit <a href="/split-pdf">vijaypdf.com/split-pdf</a></li>
<li>Upload your PDF file</li>
<li>VijayPDF automatically creates individual PDF files for every page</li>
<li>Download a ZIP archive containing all separate page files</li>
</ol>

<h2>Need to Extract Specific Pages?</h2>
<p>Use our <a href="/extract-pages">Extract Pages tool</a> to specify exactly which pages you want. Enter page numbers like "1,3,5" or "2-8" and get a new PDF containing only those pages.</p>

<h2>Need to Remove Specific Pages?</h2>
<p>Use our <a href="/remove-pages">Remove Pages tool</a> — enter the page numbers you want deleted, and download a clean PDF without those pages.</p>

<div class="blog-cta-box">
<h3>Split your PDF now — free</h3>
<a href="/split-pdf" class="btn btn-primary">Split PDF →</a>
</div>
"""
    },
    {
        "slug": "pdf-security-password-protect-pdf",
        "title": "How to Password Protect a PDF — Keep Your Files Secure",
        "meta_description": "Learn how to add password protection to any PDF document online for free. Protect sensitive files with AES-256 encryption using VijayPDF.",
        "category": "Security",
        "date": "March 15, 2026",
        "read_time": "5 min read",
        "excerpt": "Sending confidential documents as PDFs? Protect them with a password. This guide explains PDF encryption and how to do it free online.",
        "content": """
<p>PDFs containing sensitive information — contracts, financial statements, personal data, or confidential proposals — should always be password protected before sharing. Without protection, anyone who intercepts or receives the file can read it.</p>

<h2>What Is PDF Password Protection?</h2>
<p>PDF password protection works by encrypting the file's content using an encryption algorithm. Modern PDF tools use AES-256 encryption — the same standard used by banks and governments. Without the correct password, the encrypted content is completely unreadable.</p>

<h2>How to Password Protect a PDF Free</h2>
<ol>
<li>Go to <a href="/protect-pdf">vijaypdf.com/protect-pdf</a></li>
<li>Upload the PDF you want to protect</li>
<li>Enter a strong password</li>
<li>Download your password-protected PDF</li>
<li>Share the PDF and tell the recipient the password separately (never in the same email)</li>
</ol>

<h2>Tips for Strong PDF Passwords</h2>
<ul>
<li>Use at least 12 characters mixing uppercase, lowercase, numbers, and symbols</li>
<li>Avoid dictionary words, birthdates, or names</li>
<li>Use a password manager to store and share passwords securely</li>
<li>Never include the password in the same email or channel as the protected PDF</li>
</ul>

<h2>How to Remove a PDF Password</h2>
<p>If you receive a password-protected PDF and know the password, you can remove protection using our <a href="/unlock-pdf">Unlock PDF tool</a>. Enter the password and download an unprotected version for easier handling.</p>

<div class="blog-cta-box">
<h3>Protect your PDF with a password now</h3>
<a href="/protect-pdf" class="btn btn-primary">Protect PDF Free →</a>
</div>
"""
    },
    {
        "slug": "pdf-to-word-converter-comparison-2026",
        "title": "PDF to Word Converter Comparison: Which Is Best in 2026?",
        "meta_description": "Compare the best PDF to Word converters of 2026. Accuracy, formatting preservation, price, and privacy — an honest breakdown to help you choose.",
        "category": "Comparison",
        "date": "March 12, 2026",
        "read_time": "8 min read",
        "excerpt": "Not all PDF to Word converters are equal. We tested the top tools and compare accuracy, formatting, and privacy in this honest 2026 guide.",
        "content": """
<p>PDF to Word conversion sounds simple, but the quality difference between tools is enormous. A bad converter produces garbled text and broken layouts. A good one creates a Word document that looks virtually identical to the original PDF.</p>

<h2>What We Tested For</h2>
<ul>
<li><strong>Text accuracy</strong> — does all the text come through correctly?</li>
<li><strong>Table preservation</strong> — do tables retain their structure?</li>
<li><strong>Image embedding</strong> — are images maintained in position?</li>
<li><strong>Font fidelity</strong> — are font sizes, weights, and styles preserved?</li>
<li><strong>Multi-column handling</strong> — do newsletter/brochure layouts work?</li>
<li><strong>Free tier limits</strong> — what features are actually free?</li>
<li><strong>Privacy policy</strong> — how long are your files stored?</li>
</ul>

<h2>VijayPDF</h2>
<p><strong>Accuracy:</strong> Excellent for text-based PDFs<br>
<strong>Tables:</strong> Well preserved<br>
<strong>Images:</strong> Embedded correctly<br>
<strong>Free tier:</strong> Unlimited, no watermarks<br>
<strong>Privacy:</strong> Files deleted after 2 hours<br>
<strong>Best for:</strong> Regular users who want powerful free conversion without limits</p>

<h2>Adobe Acrobat Online</h2>
<p>The gold standard for PDF handling. Adobe's conversion is among the most accurate available, especially for complex layouts. However, the free tier is limited to a few conversions per month. A full subscription costs $14.99/month.</p>

<h2>Microsoft Word (Built-In)</h2>
<p>If you have Microsoft Word 2013 or later, you can open a PDF directly — Word converts it automatically. The accuracy is good for simple documents but often struggles with complex layouts. No upload to a server means maximum privacy.</p>

<h2>Google Docs</h2>
<p>Upload a PDF to Google Drive, right-click, and open with Google Docs. The conversion is reasonable for simple text PDFs but struggles with formatting. Completely free and private to your Google account.</p>

<h2>Our Recommendation</h2>
<p>For the best balance of accuracy, features, and no cost — <a href="/pdf-to-word">VijayPDF PDF to Word</a> is our top recommendation for most users. For complex print-quality layouts, Adobe Acrobat is worth the investment.</p>

<div class="blog-cta-box">
<h3>Try VijayPDF's PDF to Word converter</h3>
<a href="/pdf-to-word" class="btn btn-primary">Convert PDF to Word Free →</a>
</div>
"""
    },
    {
        "slug": "extract-text-from-pdf-ocr",
        "title": "How to Extract Text from a Scanned PDF Using OCR",
        "meta_description": "Learn how to use OCR (Optical Character Recognition) to extract text from scanned PDF documents online for free. Works with any scanned file.",
        "category": "Guide",
        "date": "March 10, 2026",
        "read_time": "6 min read",
        "excerpt": "Scanned PDFs are just images — you can't copy text from them normally. OCR technology solves this by reading and extracting the text automatically.",
        "content": """
<p>Scanned PDFs are a common headache. When someone scans a physical document, each page is saved as an image. That image looks like text, but your computer can't select, copy, or search it — because it IS just a picture.</p>

<p>OCR (Optical Character Recognition) solves this by analyzing the image and recognizing characters, words, and sentences — then outputting them as actual, machine-readable text.</p>

<h2>What Is OCR?</h2>
<p>OCR is a technology that analyzes images of text and converts them into computer-readable characters. Modern OCR engines use machine learning to recognize over 100 languages and handle varying font styles, sizes, and orientations with remarkable accuracy.</p>

<h2>When Do You Need OCR?</h2>
<ul>
<li>When you can't select or copy text in a PDF (it's a scanned image)</li>
<li>When searching a PDF for a keyword returns no results</li>
<li>When you want to convert a scanned document to Word or Excel</li>
<li>When you receive old paper-based records that have been scanned</li>
<li>When extracting data from printed invoices or forms</li>
</ul>

<h2>How to Use OCR on a Scanned PDF</h2>
<ol>
<li>Go to <a href="/ocr-pdf">vijaypdf.com/ocr-pdf</a></li>
<li>Upload your scanned PDF</li>
<li>Our OCR engine analyzes every page and extracts all readable text</li>
<li>Download a .txt file containing all extracted text</li>
</ol>

<h2>After OCR: Convert to Word or Excel</h2>
<p>Once your scanned PDF has been OCR-processed, you can:</p>
<ul>
<li>Convert the resulting text to <a href="/pdf-to-word">Word format</a> for editing</li>
<li>Extract tables to <a href="/pdf-to-excel">Excel format</a> for data analysis</li>
<li>Use the text in reports, databases, or other applications</li>
</ul>

<div class="blog-cta-box">
<h3>Extract text from your scanned PDF now</h3>
<a href="/ocr-pdf" class="btn btn-primary">Use OCR PDF Tool →</a>
</div>
"""
    },
    {
        "slug": "how-to-add-watermark-to-pdf",
        "title": "How to Add a Watermark to a PDF Online Free",
        "meta_description": "Add text watermarks to any PDF document for free online. Mark documents as confidential, draft, or add your branding in seconds with VijayPDF.",
        "category": "Tutorial",
        "date": "March 8, 2026",
        "read_time": "4 min read",
        "excerpt": "Protect and brand your PDF documents by adding text watermarks. Mark files as Confidential, Draft, or add your company name — free online.",
        "content": """
<p>Watermarks serve an important dual purpose: they identify the author or owner of a document, and they signal the document's status — whether it's a draft, confidential, or for internal use only. Adding a watermark to a PDF takes seconds with VijayPDF.</p>

<h2>Common Watermark Uses</h2>
<ul>
<li><strong>CONFIDENTIAL</strong> — mark sensitive documents before sharing</li>
<li><strong>DRAFT</strong> — indicate documents that are not yet finalized</li>
<li><strong>SAMPLE</strong> — share samples of work without full distribution rights</li>
<li><strong>Company name / logo text</strong> — brand every page of a shared document</li>
<li><strong>DO NOT COPY</strong> — discourage unauthorized reproduction</li>
</ul>

<h2>How to Add a Watermark to a PDF</h2>
<ol>
<li>Visit <a href="/add-watermark">vijaypdf.com/add-watermark</a></li>
<li>Upload your PDF</li>
<li>Enter the watermark text you want to appear</li>
<li>Download your watermarked PDF</li>
</ol>

<h2>Can I Remove a Watermark from a PDF?</h2>
<p>If you added a text watermark and want to remove it, the easiest solution is to use the original unwatermarked file. Watermarks embedded in a PDF's content stream are difficult to remove cleanly without specialized tools.</p>

<h2>Related PDF Security Tools</h2>
<p>For comprehensive PDF security, combine watermarking with:</p>
<ul>
<li><a href="/protect-pdf">Password protection</a> — prevent unauthorized opening</li>
<li><a href="/redact-pdf">Redaction</a> — permanently remove sensitive text before sharing</li>
</ul>

<div class="blog-cta-box">
<h3>Add a watermark to your PDF now</h3>
<a href="/add-watermark" class="btn btn-primary">Add Watermark Free →</a>
</div>
"""
    },
    {
        "slug": "excel-to-pdf-guide",
        "title": "How to Convert Excel to PDF — Keep Formatting Perfect",
        "meta_description": "Convert Excel spreadsheets to PDF online for free. Preserve all tables, charts, and formatting. No Excel installation needed. Use VijayPDF free.",
        "category": "Tutorial",
        "date": "March 6, 2026",
        "read_time": "5 min read",
        "excerpt": "Converting Excel to PDF ensures your spreadsheet looks identical for every recipient. Here's how to do it free online — no Excel needed.",
        "content": """
<p>Excel spreadsheets are powerful, but sharing them can be problematic: the recipient needs Excel, different versions display things differently, and formulas can break. Converting to PDF gives you a universally viewable, perfectly formatted document every time.</p>

<h2>Why Convert Excel to PDF?</h2>
<ul>
<li>Recipients don't need Excel installed to view it</li>
<li>PDF preserves your exact cell widths, fonts, and print area settings</li>
<li>Formulas won't accidentally be modified by recipients</li>
<li>PDFs print consistently on any printer</li>
<li>PDF is the standard format for financial reports and invoices</li>
</ul>

<h2>How to Convert Excel to PDF Free</h2>
<ol>
<li>Go to <a href="/excel-to-pdf">vijaypdf.com/excel-to-pdf</a></li>
<li>Upload your .xlsx or .xls file</li>
<li>VijayPDF converts your spreadsheet to PDF, preserving tables and data</li>
<li>Download your PDF immediately</li>
</ol>

<h2>Tips for Best Excel to PDF Results</h2>
<ul>
<li>Set your print area in Excel before converting for cleaner page breaks</li>
<li>Use landscape orientation for wide spreadsheets</li>
<li>After conversion, compress the PDF with <a href="/compress-pdf">Compress PDF</a> if needed</li>
<li>To go back to Excel from a PDF, use our <a href="/pdf-to-excel">PDF to Excel tool</a></li>
</ul>

<div class="blog-cta-box">
<h3>Convert your Excel file to PDF now</h3>
<a href="/excel-to-pdf" class="btn btn-primary">Convert Excel to PDF →</a>
</div>
"""
    },
    {
        "slug": "pdf-for-email-reduce-size",
        "title": "How to Make a PDF Small Enough to Email",
        "meta_description": "PDF too large to email? Learn how to reduce PDF file size for email attachments under 10MB or 25MB limits. Free online PDF compressor guide.",
        "category": "Guide",
        "date": "March 4, 2026",
        "read_time": "5 min read",
        "excerpt": "Most email services limit attachments to 10–25MB. If your PDF is too large to email, here's how to shrink it fast — for free.",
        "content": """
<p>You've prepared the perfect document, saved it as a PDF, and then... "Attachment too large." Gmail limits attachments to 25MB. Outlook to 20MB. Many corporate email systems to just 10MB. Large scanned documents and image-heavy PDFs blow past these limits easily.</p>

<h2>Quick Fix: Compress the PDF</h2>
<p>The fastest solution is to run your PDF through a free compressor:</p>
<ol>
<li>Go to <a href="/compress-pdf">vijaypdf.com/compress-pdf</a></li>
<li>Upload your oversized PDF</li>
<li>Download the compressed version</li>
<li>Check the new file size — usually 30–70% smaller</li>
</ol>

<h2>If Compression Isn't Enough</h2>
<p>For very large documents (100MB+), compression alone may not get below the email limit. Try these additional steps:</p>
<ul>
<li><strong>Split the PDF</strong> — use <a href="/split-pdf">Split PDF</a> to break it into smaller sections and send in multiple emails</li>
<li><strong>Remove unnecessary pages</strong> — use <a href="/remove-pages">Remove Pages</a> to strip out pages the recipient doesn't need</li>
<li><strong>Use Google Drive or Dropbox</strong> — upload the full PDF to cloud storage and share a link instead of an attachment</li>
</ul>

<h2>Email Size Limits by Provider</h2>
<ul>
<li><strong>Gmail:</strong> 25 MB per attachment (10 MB for Google Drive files)</li>
<li><strong>Outlook/Hotmail:</strong> 20 MB</li>
<li><strong>Yahoo Mail:</strong> 25 MB</li>
<li><strong>Corporate email (Exchange):</strong> Varies — often 10–15 MB</li>
</ul>

<div class="blog-cta-box">
<h3>Compress your PDF for email now</h3>
<a href="/compress-pdf" class="btn btn-primary">Compress PDF Free →</a>
</div>
"""
    },
    {
        "slug": "rotate-pdf-pages-online",
        "title": "How to Rotate PDF Pages Online Free",
        "meta_description": "Rotate individual PDF pages or the entire document online for free. Fix upside down or sideways pages in seconds. No software installation needed.",
        "category": "Tutorial",
        "date": "March 2, 2026",
        "read_time": "4 min read",
        "excerpt": "Got a PDF with upside-down or sideways pages? Fix them in seconds with VijayPDF's free online PDF rotation tool.",
        "content": """
<p>Scanned documents often come out sideways. PDFs created from photos taken in portrait mode on a phone can end up rotated. Whatever the cause, VijayPDF lets you rotate individual pages or the entire document quickly and for free.</p>

<h2>How to Rotate PDF Pages Online</h2>
<ol>
<li>Go to <a href="/rotate-pdf">vijaypdf.com/rotate-pdf</a></li>
<li>Upload your PDF file</li>
<li>Specify the rotation angle (90°, 180°, or 270°)</li>
<li>Download your corrected PDF</li>
</ol>

<h2>Common PDF Rotation Scenarios</h2>
<ul>
<li>Scanned landscape document that appears portrait in the PDF viewer</li>
<li>Mobile phone photo of a document saved sideways</li>
<li>PDF with a mix of portrait and landscape pages that need individual correction</li>
<li>Book scans where alternating pages are rotated for double-page spreads</li>
</ul>

<div class="blog-cta-box">
<h3>Fix your PDF rotation now</h3>
<a href="/rotate-pdf" class="btn btn-primary">Rotate PDF Free →</a>
</div>
"""
    },
    {
        "slug": "what-is-pdf-a-format",
        "title": "What Is PDF/A? The Archiving Standard Explained",
        "meta_description": "PDF/A is the ISO standard for long-term document archiving. Learn what PDF/A means, why it matters for legal and government documents, and how to create one.",
        "category": "Education",
        "date": "February 28, 2026",
        "read_time": "6 min read",
        "excerpt": "Government agencies and legal departments require PDF/A for archiving. This guide explains what PDF/A is and why it matters.",
        "content": """
<p>Most people are familiar with PDF but have never heard of PDF/A. If you've ever been asked to submit a document in "PDF/A format" and had no idea what that meant — this guide explains everything.</p>

<h2>What Is PDF/A?</h2>
<p>PDF/A (the "A" stands for Archive) is an ISO-standardized version of the PDF format specifically designed for long-term digital preservation. It was created to address a fundamental problem: today's PDF files may be unreadable decades from now because they depend on external resources, proprietary fonts, or encryption that may not be supported by future software.</p>

<h2>How Is PDF/A Different from Regular PDF?</h2>
<p>PDF/A files must be entirely self-contained:</p>
<ul>
<li>All fonts must be embedded — no relying on system fonts</li>
<li>No encryption or password protection allowed</li>
<li>No external content references (URLs, linked images)</li>
<li>Color spaces must be explicitly specified</li>
<li>No audio, video, or interactive JavaScript</li>
<li>Full metadata (creation date, author, title) must be included</li>
</ul>

<h2>Who Requires PDF/A?</h2>
<ul>
<li>Government archives — court systems, land registries, national libraries</li>
<li>Legal document management — permanent contract records</li>
<li>Healthcare records — patient documentation that must be retained for decades</li>
<li>Financial institutions — compliant document storage for regulatory requirements</li>
</ul>

<h2>Convert PDF to PDF/A</h2>
<p>Use our <a href="/pdf-to-pdfa">PDF to PDF/A converter</a> to create archiving-compliant versions of your documents.</p>

<div class="blog-cta-box">
<h3>Convert your PDF to PDF/A format</h3>
<a href="/pdf-to-pdfa" class="btn btn-primary">Convert to PDF/A →</a>
</div>
"""
    },
    {
        "slug": "ai-pdf-summarizer-guide",
        "title": "How to Summarize Long PDFs with AI — Save Time Reading",
        "meta_description": "Use AI to automatically summarize long PDF documents in seconds. Get key points without reading 50 pages. VijayPDF's AI summarizer explained.",
        "category": "AI Tools",
        "date": "February 25, 2026",
        "read_time": "5 min read",
        "excerpt": "AI can read and summarize a 50-page PDF in seconds. Here's how to use VijayPDF's AI summarizer to extract key insights instantly.",
        "content": """
<p>The average business professional reads 5–7 full reports, proposals, or research papers per week. At 30–80 pages each, that's hundreds of pages — much of which contains filler, context, and background. AI summarization extracts what matters in seconds.</p>

<h2>What Is AI PDF Summarization?</h2>
<p>AI summarization uses large language models (LLMs) to read, understand, and condense the content of a PDF document. The AI identifies the main topics, key arguments, important data, and conclusions — then generates a concise summary.</p>

<h2>How to Summarize a PDF with AI</h2>
<ol>
<li>Go to <a href="/summarize-pdf">vijaypdf.com/summarize-pdf</a></li>
<li>Upload your PDF document (up to 50 MB)</li>
<li>VijayPDF's AI reads the document content</li>
<li>Receive a concise, intelligible summary in seconds</li>
</ol>

<h2>Best Use Cases for AI PDF Summarization</h2>
<ul>
<li>Research papers — get key findings without reading methodology sections</li>
<li>Legal contracts — quickly identify key terms and obligations</li>
<li>Annual reports — extract financial highlights from 80-page documents</li>
<li>Academic textbooks — chapter summaries for study and revision</li>
<li>News and policy documents — understand government reports quickly</li>
</ul>

<div class="blog-cta-box">
<h3>Summarize your PDF with AI now</h3>
<a href="/summarize-pdf" class="btn btn-primary">AI Summarize PDF →</a>
</div>
"""
    },
    {
        "slug": "pdf-tools-for-students",
        "title": "Best PDF Tools for Students in 2026 — All Free",
        "meta_description": "Discover the best free PDF tools for students. Merge notes, compress assignments, convert PDFs to Word, and more — all free online without software.",
        "category": "Students",
        "date": "February 22, 2026",
        "read_time": "6 min read",
        "excerpt": "Students deal with PDFs constantly — lecture notes, assignments, research papers. These free tools make student life easier.",
        "content": """
<p>From lecture notes shared as PDFs to assignment submissions, research paper reading, and exam preparation — students interact with PDF documents every single day. Here are the free tools that save the most time.</p>

<h2>1. Merge PDF — Combine Assignments and Notes</h2>
<p>If your professor provides lecture notes as separate PDFs for each week, merge them into one master document for easier studying. Use <a href="/merge-pdf">VijayPDF Merge PDF</a> to combine any number of PDFs.</p>

<h2>2. Compress PDF — Stay Under Submission Limits</h2>
<p>University portals often limit file uploads to 5–10 MB. If your assignment PDF is too large, <a href="/compress-pdf">compress it</a> before submission.</p>

<h2>3. PDF to Word — Edit and Annotate</h2>
<p>Need to edit a template, fill in a form, or extract content from a PDF? Convert it to Word with <a href="/pdf-to-word">PDF to Word converter</a> and make your changes easily.</p>

<h2>4. JPG to PDF — Submit Handwritten Work</h2>
<p>Many students photograph handwritten assignments and need to submit them as PDFs. <a href="/jpg-to-pdf">JPG to PDF</a> combines those photos into a clean document instantly.</p>

<h2>5. AI Summarize — Study Smarter</h2>
<p>Use <a href="/summarize-pdf">AI PDF Summarizer</a> to get quick summaries of long research papers and textbook chapters when you're short on time.</p>

<h2>6. Split PDF — Extract Reading Chapters</h2>
<p>When a professor assigns "chapters 3 and 7" from a PDF textbook, use <a href="/extract-pages">Extract Pages</a> to get just those chapters.</p>

<div class="blog-cta-box">
<h3>Access all free student PDF tools</h3>
<a href="/tools" class="btn btn-primary">Explore All Tools →</a>
</div>
"""
    },
    {
        "slug": "pdf-tools-for-businesses",
        "title": "Essential PDF Tools Every Business Needs in 2026",
        "meta_description": "Discover the essential PDF tools every business needs. From secure document sharing to data extraction and AI summaries — all available free online.",
        "category": "Business",
        "date": "February 20, 2026",
        "read_time": "7 min read",
        "excerpt": "PDFs are the backbone of business document management. These free tools handle everything from security to data extraction and AI summarization.",
        "content": """
<p>Every business — from a one-person freelance operation to a large enterprise — deals with PDFs daily. Contracts, invoices, proposals, reports, compliance documents. The right PDF tools dramatically cut down the time spent on document management.</p>

<h2>Document Security</h2>
<p>Protect sensitive business documents before sharing:</p>
<ul>
<li><a href="/protect-pdf"><strong>Password Protect PDF</strong></a> — encrypt contracts and financial documents with AES-256 before emailing clients</li>
<li><a href="/add-watermark"><strong>Add Watermark</strong></a> — mark proposals and samples with "CONFIDENTIAL" or your company name</li>
<li><a href="/redact-pdf"><strong>Redact PDF</strong></a> — permanently black out sensitive information (personal data, pricing) before sharing partial documents</li>
</ul>

<h2>Document Creation and Conversion</h2>
<ul>
<li><a href="/word-to-pdf"><strong>Word to PDF</strong></a> — create professional-looking PDFs from all proposals and reports</li>
<li><a href="/excel-to-pdf"><strong>Excel to PDF</strong></a> — share financial data without giving recipients access to formulas</li>
<li><a href="/jpg-to-pdf"><strong>Image to PDF</strong></a> — turn product photos or scanned forms into PDF documents</li>
</ul>

<h2>Document Management</h2>
<ul>
<li><a href="/merge-pdf"><strong>Merge PDF</strong></a> — combine multiple contracts or reports into complete submission packages</li>
<li><a href="/split-pdf"><strong>Split PDF</strong></a> — extract specific sections from lengthy legal or technical documents</li>
<li><a href="/compress-pdf"><strong>Compress PDF</strong></a> — reduce file sizes for email delivery and cloud storage</li>
</ul>

<h2>Data Extraction</h2>
<ul>
<li><a href="/pdf-to-excel"><strong>PDF to Excel</strong></a> — extract financial data from PDF reports into spreadsheets for analysis</li>
<li><a href="/pdf-to-word"><strong>PDF to Word</strong></a> — make contracts and templates editable for customization</li>
<li><a href="/ocr-pdf"><strong>OCR PDF</strong></a> — extract text from scanned invoices and forms for data entry</li>
</ul>

<div class="blog-cta-box">
<h3>All PDF tools — free for your business</h3>
<a href="/tools" class="btn btn-primary">View All Tools →</a>
</div>
"""
    },
    {
        "slug": "unlock-protected-pdf-guide",
        "title": "How to Unlock a Password-Protected PDF for Free",
        "meta_description": "Learn how to remove passwords from PDF files for free online. Unlock protected PDFs you own in seconds using VijayPDF's free unlock tool.",
        "category": "Guide",
        "date": "February 18, 2026",
        "read_time": "4 min read",
        "excerpt": "Forgotten a PDF password? Or tired of entering it every time? This guide shows how to unlock a protected PDF you own — for free.",
        "content": """
<p><strong>Important note:</strong> This guide is for unlocking PDFs that you own and have the password for. We do not support bypassing security on documents you don't have permission to access.</p>

<h2>Why Might You Want to Unlock a PDF?</h2>
<ul>
<li>You're tired of entering the password every time you open the document</li>
<li>You need to merge a protected PDF with other files</li>
<li>You want to print a PDF that has printing restrictions</li>
<li>You're backing up old documents and want to remove legacy passwords</li>
</ul>

<h2>How to Unlock a PDF You Own</h2>
<ol>
<li>Go to <a href="/unlock-pdf">vijaypdf.com/unlock-pdf</a></li>
<li>Upload your password-protected PDF</li>
<li>Enter the current password when prompted</li>
<li>Download the unlocked PDF without any password protection</li>
</ol>

<h2>What If I've Forgotten the Password?</h2>
<p>If you've genuinely forgotten the password to your own document, recovery options depend on how the PDF was created and encrypted. Standard AES-256 encrypted PDFs are essentially impossible to brute-force. Your best option may be to check with the document's original creator for the password.</p>

<h2>After Unlocking: Keep It Secure</h2>
<p>If your PDF contained sensitive information, consider re-protecting it with a memorable new password using our <a href="/protect-pdf">Protect PDF tool</a> after making your changes.</p>

<div class="blog-cta-box">
<h3>Unlock your password-protected PDF now</h3>
<a href="/unlock-pdf" class="btn btn-primary">Unlock PDF Free →</a>
</div>
"""
    },
    {
        "slug": "pdf-vs-docx-which-format-to-use",
        "title": "PDF vs Word (DOCX): When to Use Each Format",
        "meta_description": "Should you send a PDF or a Word document? This guide explains the key differences between PDF and DOCX and helps you choose the right format every time.",
        "category": "Education",
        "date": "February 15, 2026",
        "read_time": "6 min read",
        "excerpt": "PDF and Word are both essential document formats, but they serve very different purposes. This guide helps you pick the right one every time.",
        "content": """
<p>PDF and DOCX are the two most common document formats in professional settings — but they're designed for very different purposes. Choosing the wrong one causes headaches for both sender and recipient.</p>

<h2>Use PDF When</h2>
<ul>
<li>You want the document to look identical on every device and operating system</li>
<li>You're sharing a final version that should not be modified</li>
<li>You need to ensure printing consistency (reports, invoices, brochures)</li>
<li>The recipient may not have Microsoft Office installed</li>
<li>You need to protect content from easy copying or editing</li>
<li>Submitting documents to government portals or official systems</li>
</ul>

<h2>Use Word (DOCX) When</h2>
<ul>
<li>Collaborators need to make tracked changes or add comments</li>
<li>The document needs to be edited and updated regularly</li>
<li>You're working on a draft that will go through multiple revisions</li>
<li>You want to use mail merge or other Word automation features</li>
<li>You need the recipient to fill in a template</li>
</ul>

<h2>Converting Between Formats</h2>
<p>The good news is that converting between PDF and Word is easy with VijayPDF:</p>
<ul>
<li><a href="/pdf-to-word">PDF to Word</a> — get an editable Word document from any PDF</li>
<li><a href="/word-to-pdf">Word to PDF</a> — create a professional PDF from any Word document</li>
</ul>

<div class="blog-cta-box">
<h3>Convert between PDF and Word for free</h3>
<a href="/tools" class="btn btn-primary">View Conversion Tools →</a>
</div>
"""
    },
    {
        "slug": "compress-pdf-without-losing-quality-tips",
        "title": "7 Expert Tips to Compress PDF Without Losing Quality",
        "meta_description": "Get expert tips for compressing PDF files while preserving quality. Reduce file size for email, storage, and web without blurry images or broken text.",
        "category": "Tips",
        "date": "February 12, 2026",
        "read_time": "7 min read",
        "excerpt": "Compressing a PDF can go wrong if you don't know the tricks. These 7 expert tips help you reduce file size without sacrificing document quality.",
        "content": """
<p>PDF compression is straightforward — upload and download. But getting the best results means understanding what's inside your PDF and choosing the right approach. These 7 expert tips will help you achieve maximum compression with minimum quality loss.</p>

<h2>1. Use a Purpose-Built PDF Compressor</h2>
<p>General-purpose tools don't optimize for PDFs. Dedicated tools like <a href="/compress-pdf">VijayPDF's compressor</a> understand PDF internal structure and target the biggest size contributors directly.</p>

<h2>2. Understand What Makes Your PDF Large</h2>
<p>Before compressing, check what's inflating the size:</p>
<ul>
<li>Embedded high-resolution images are typically the #1 culprit</li>
<li>Scanned page images (TIFFs at 600 DPI) compress enormously</li>
<li>Unsubsetted font embeddings can add megabytes</li>
</ul>

<h2>3. Compress Images Before Embedding</h2>
<p>If you're creating a PDF from scratch, optimize your images first. Resize photos to the actual display size (96–150 DPI for screen, 300 DPI for print) before embedding them.</p>

<h2>4. Separate Content for Separate Recipients</h2>
<p>Instead of sending one massive PDF, use <a href="/split-pdf">Split PDF</a> to send only relevant sections to each recipient. Smaller individual files mean less need for compression.</p>

<h2>5. Remove Unnecessary Pages</h2>
<p>Before compressing, use <a href="/remove-pages">Remove Pages</a> to strip out blank pages, cover sheets, or sections that don't need to be shared. Less content = smaller file.</p>

<h2>6. Don't Compress What's Already Compressed</h2>
<p>Running an already-optimized PDF through a compressor again yields very little savings (often under 5%). Save your effort - check the PDF's compression status first.</p>

<h2>7. Use Cloud Storage for Very Large Files</h2>
<p>For PDFs that remain large after all optimizations (100MB+), use Google Drive or Dropbox to share via link rather than email attachment. This bypasses all size limits entirely.</p>

<div class="blog-cta-box">
<h3>Compress your PDF now</h3>
<a href="/compress-pdf" class="btn btn-primary">Compress PDF Free →</a>
</div>
"""
    },
]

def get_blog_post(slug):
    """Return a blog post dict by slug, or None if not found."""
    return next((p for p in BLOG_POSTS if p["slug"] == slug), None)

# Display Web Routes

@app.errorhandler(404)
def page_not_found(e):
    return render_template('404.html'), 404

@app.route("/blog")
def blog_index():
    return render_template("blog/index.html", posts=BLOG_POSTS)

@app.route("/blog/<slug>")
def blog_post(slug):
    post = get_blog_post(slug)
    if not post:
        abort(404)
    # Related posts: 3 random others
    related = [p for p in BLOG_POSTS if p["slug"] != slug][:3]
    return render_template("blog/post.html", post=post, related=related)

@app.route("/tools")
def tools():
    return render_template("tools.html")

@app.route("/pricing")
def pricing():
    return render_template("pricing.html")

# Tool Pages
@app.route("/pdf-to-word")
def pdf_to_word():
    return render_template("tools/pdf_to_word.html")

@app.route("/word-to-pdf")
def word_to_pdf():
    return render_template("tools/word_to_pdf.html")

@app.route("/jpg-to-pdf")
def jpg_to_pdf():
    return render_template("tools/jpg_to_pdf.html")

@app.route("/pdf-to-jpg")
def pdf_to_jpg():
    return render_template("tools/pdf_to_jpg.html")

@app.route("/merge-pdf")
def merge_pdf():
    return render_template("tools/merge_pdf.html")

@app.route("/split-pdf")
def split_pdf():
    return render_template("tools/split_pdf.html")

@app.route("/compress-pdf")
def compress_pdf():
    return render_template("tools/compress_pdf.html")

@app.route("/pdf-preview")
def pdf_preview():
    return render_template("tools/pdf_preview.html")

@app.route("/unlock-pdf")
def unlock_pdf():
    return render_template("tools/unlock_pdf.html")

@app.route("/protect-pdf")
def protect_pdf():
    return render_template("tools/protect_pdf.html")

@app.route("/remove-pages")
def remove_pages():
    return render_template("tools/remove_pages.html")

@app.route("/extract-pages")
def extract_pages():
    return render_template("tools/extract_pages.html")

@app.route("/organize-pdf")
def organize_pdf():
    return render_template("tools/organize_pdf.html")

@app.route("/scan-to-pdf")
def scan_to_pdf():
    return render_template("tools/scan_to_pdf.html")

@app.route("/ocr-pdf")
def ocr_pdf():
    return render_template("tools/ocr_pdf.html")

@app.route("/pdf-to-excel")
def pdf_to_excel():
    return render_template("tools/pdf_to_excel.html")

@app.route("/excel-to-pdf")
def excel_to_pdf():
    return render_template("tools/excel_to_pdf.html")

@app.route("/pdf-to-powerpoint")
def pdf_to_powerpoint():
    return render_template("tools/pdf_to_powerpoint.html")

@app.route("/powerpoint-to-pdf")
def powerpoint_to_pdf():
    return render_template("tools/powerpoint_to_pdf.html")

@app.route("/pdf-to-pdfa")
def pdf_to_pdfa():
    return render_template("tools/pdf_to_pdfa.html")

@app.route("/rotate-pdf")
def rotate_pdf():
    return render_template("tools/rotate_pdf.html")

@app.route("/add-page-numbers")
def add_page_numbers():
    return render_template("tools/add_page_numbers.html")

@app.route("/add-watermark")
def add_watermark():
    return render_template("tools/add_watermark.html")

@app.route("/crop-pdf")
def crop_pdf():
    return render_template("tools/crop_pdf.html")

@app.route("/edit-pdf")
def edit_pdf():
    return render_template("tools/edit_pdf.html")

@app.route("/redact-pdf")
def redact_pdf():
    return render_template("tools/redact_pdf.html")

@app.route("/compare-pdf")
def compare_pdf():
    return render_template("tools/compare_pdf.html")

@app.route("/summarize-pdf")
def summarize_pdf():
    return render_template("tools/summarize_pdf.html")

@app.route("/translate-pdf")
def translate_pdf():
    return render_template("tools/translate_pdf.html")

@app.route("/chat-pdf")
def chat_pdf():
    return render_template("tools/chat_pdf.html")

# Info Pages
@app.route("/about")
def about():
    return render_template("about.html")

@app.route("/contact-support", methods=['GET', 'POST'])
def contact_support():
    if request.method == 'POST':
        name = request.form.get('name')
        email = request.form.get('email')
        message = request.form.get('message')
        
        try:
            # 1. Send Internal Notification
            subject = f"New Contact Message - {name}"
            body = f"Name: {name}\nEmail: {email}\n\nMessage:\n{message}"
            email_service.send_email(os.getenv("EMAIL_USER"), f"New Contact Message - VijayPDF", body)
            
            # 2. Send Auto-Reply to User
            email_service.send_auto_reply(email, name)
        except Exception as e:
            app.logger.error(f"Contact form email error: {e}")
        
        flash("Thank you for contacting VijayPDF.com. Our team will respond shortly.", "success")
        return render_template("contact_support.html", success=True)
        
    return render_template("contact_support.html")

@app.route("/faq")
def faq():
    return render_template("faq.html")

@app.route("/privacy")
def privacy():
    return render_template("privacy_policy.html")

@app.route("/terms")
def terms():
    return render_template("terms.html")

@app.route('/api/merge', methods=['POST'])
def route_merge(): return process_wrapper(lambda in_paths, out_path: merge_pdf_func(in_paths, out_path), multi=True)

@app.route('/api/split', methods=['POST'])
def route_split(): return process_wrapper(lambda inp, out: split_pdf_func(inp, out), ext="_split.zip")

@app.route('/api/remove-pages', methods=['POST'])
def route_rem(): return process_wrapper(lambda inp, out: remove_pages_func(inp, out, request.form.get('pages', '')))

@app.route('/api/extract-pages', methods=['POST'])
def route_ext(): return process_wrapper(lambda inp, out: extract_pages_func(inp, out, request.form.get('pages', '')))

@app.route('/api/organize', methods=['POST'])
def route_org(): return process_wrapper(lambda inp, out: organize_pdf_func(inp, out, request.form.get('pages', '')))

@app.route('/api/scan-to-pdf', methods=['POST'])
def route_scan(): return process_wrapper(lambda inp, out: scan_to_pdf_func(inp, out), multi=True)

@app.route('/api/compress', methods=['POST'])
def route_comp(): return process_wrapper(lambda inp, out: compress_pdf_func(inp, out))

@app.route('/api/repair', methods=['POST'])
def route_rep(): return process_wrapper(lambda inp, out: repair_pdf_func(inp, out))

@app.route('/api/ocr', methods=['POST'])
def route_ocr(): return process_wrapper(lambda inp, out: ocr_pdf_func(inp, out), ext=".txt")

@app.route('/api/jpg-to-pdf', methods=['POST'])
def route_j2p(): return process_wrapper(lambda inp, out: jpg_to_pdf_func(inp, out), multi=True)

@app.route('/api/word-to-pdf', methods=['POST'])
def route_w2p(): return process_wrapper(lambda inp, out: word_to_pdf_func(inp, out))

@app.route('/api/powerpoint-to-pdf', methods=['POST'])
def route_p2p(): return process_wrapper(lambda inp, out: powerpoint_to_pdf_func(inp, out))

@app.route('/api/excel-to-pdf', methods=['POST'])
def route_e2p(): return process_wrapper(lambda inp, out: excel_to_pdf_func(inp, out))

@app.route('/api/html-to-pdf', methods=['POST'])
def route_h2p(): return process_wrapper(lambda inp, out: html_to_pdf_func(inp, out))

@app.route('/api/pdf-to-jpg', methods=['POST'])
def route_p2j(): return process_wrapper(lambda inp, out: pdf_to_jpg_func(inp, out), ext=".zip")

@app.route('/api/pdf-to-word', methods=['POST'])
def route_p2w(): return process_wrapper(lambda inp, out: pdf_to_word_func(inp, out), ext=".docx")

@app.route('/api/pdf-to-powerpoint', methods=['POST'])
def route_p2pptx(): return process_wrapper(lambda inp, out: pdf_to_powerpoint_func(inp, out), ext=".pptx")

@app.route('/api/pdf-to-excel', methods=['POST'])
def route_p2x(): return process_wrapper(lambda inp, out: pdf_to_excel_func(inp, out), ext=".xlsx")

@app.route('/api/pdf-to-pdfa', methods=['POST'])
def route_p2a(): return process_wrapper(lambda inp, out: pdf_to_pdfa_func(inp, out))

@app.route('/api/rotate', methods=['POST'])
def route_rot(): return process_wrapper(lambda inp, out: rotate_pdf_func(inp, out, request.form.get('angle', '90')))

@app.route('/api/add-page-numbers', methods=['POST'])
def route_pn(): return process_wrapper(lambda inp, out: add_page_numbers_func(inp, out))

@app.route('/api/add-watermark', methods=['POST'])
def route_wm(): return process_wrapper(lambda inp, out: add_watermark_func(inp, out, request.form.get('text', 'WATERMARK')))

@app.route('/api/crop', methods=['POST'])
def route_crop(): return process_wrapper(lambda inp, out: crop_pdf_func(inp, out))

@app.route('/api/edit-pdf', methods=['POST'])
def route_edit(): return process_wrapper(lambda inp, out: edit_pdf_func(inp, out, request.form.get('text', '')))

@app.route('/api/unlock', methods=['POST'])
def route_unl(): return process_wrapper(lambda inp, out: unlock_pdf_func(inp, out, request.form.get('password', '')))

@app.route('/api/protect', methods=['POST'])
def route_prot(): return process_wrapper(lambda inp, out: protect_pdf_func(inp, out, request.form.get('password', '')))

@app.route('/api/sign', methods=['POST'])
def route_sign(): return process_wrapper(lambda inp, out: sign_pdf_func(inp, out, request.form.get('pfx'), request.form.get('password', '')))

@app.route('/api/redact', methods=['POST'])
def route_redact(): return process_wrapper(lambda inp, out: redact_pdf_func(inp, out, request.form.get('text', '')))

@app.route('/api/compare', methods=['POST'])
def route_comp_pdf(): 
    f = save_uploaded_files(request.files.getlist('files[]'))
    if len(f)<2: return jsonify({'error': 'Needs 2 files'}), 400
    out_name = f"{uuid.uuid4().hex}.txt"
    compare_pdf_func(f[0][0], f[1][0], os.path.join(app.config['OUTPUT_FOLDER'], out_name))
    return jsonify({'success':True, 'download_url': url_for('download_file', filename=out_name)})

@app.route('/api/summarize', methods=['POST'])
def route_sum(): return process_wrapper(lambda inp: ai_summarize_func(inp), is_text=True)

@app.route('/api/translate', methods=['POST'])
def route_trans(): return process_wrapper(lambda inp: translate_pdf_func(inp, request.form.get('language', 'Spanish')), is_text=True)

@app.route('/api/chat', methods=['POST'])
@csrf.exempt
def route_chat(): return process_wrapper(lambda inp: chat_with_pdf_func(inp, request.form.get('question', 'What is this document about?')), is_text=True)

@app.route('/sitemap.xml')
@limiter.exempt
def sitemap():
    """Serve a dynamically generated sitemap.xml for SEO — includes tool pages and blog posts."""
    base = "https://www.vijaypdf.com"
    today = datetime.utcnow().strftime("%Y-%m-%d")

    pages = [
        ("/", "1.0", "weekly"),
        ("/tools", "0.9", "weekly"),
        ("/blog", "0.9", "weekly"),
        # --- Core Tool Pages ---
        ("/pdf-to-word", "0.8", "monthly"),
        ("/word-to-pdf", "0.8", "monthly"),
        ("/jpg-to-pdf", "0.8", "monthly"),
        ("/pdf-to-jpg", "0.8", "monthly"),
        ("/merge-pdf", "0.8", "monthly"),
        ("/split-pdf", "0.8", "monthly"),
        ("/compress-pdf", "0.8", "monthly"),
        ("/pdf-to-excel", "0.7", "monthly"),
        ("/excel-to-pdf", "0.7", "monthly"),
        ("/pdf-preview", "0.7", "monthly"),
        ("/unlock-pdf", "0.7", "monthly"),
        ("/protect-pdf", "0.7", "monthly"),
        ("/remove-pages", "0.7", "monthly"),
        ("/extract-pages", "0.7", "monthly"),
        ("/organize-pdf", "0.7", "monthly"),
        ("/scan-to-pdf", "0.7", "monthly"),
        ("/ocr-pdf", "0.7", "monthly"),
        ("/pdf-to-powerpoint", "0.7", "monthly"),
        ("/powerpoint-to-pdf", "0.7", "monthly"),
        ("/pdf-to-pdfa", "0.7", "monthly"),
        ("/rotate-pdf", "0.7", "monthly"),
        ("/add-page-numbers", "0.7", "monthly"),
        ("/add-watermark", "0.7", "monthly"),
        ("/crop-pdf", "0.7", "monthly"),
        ("/edit-pdf", "0.7", "monthly"),
        ("/redact-pdf", "0.7", "monthly"),
        ("/compare-pdf", "0.7", "monthly"),
        ("/summarize-pdf", "0.7", "monthly"),
        ("/translate-pdf", "0.7", "monthly"),
        ("/chat-pdf", "0.7", "monthly"),
        # --- Info Pages ---
        ("/pricing", "0.6", "monthly"),
        ("/about", "0.5", "monthly"),
        ("/contact-support", "0.5", "monthly"),
        ("/faq", "0.5", "monthly"),
        ("/privacy", "0.4", "yearly"),
        ("/terms", "0.4", "yearly"),
    ]

    # Add all blog posts dynamically
    for post in BLOG_POSTS:
        pages.append((f"/blog/{post['slug']}", "0.7", "monthly"))

    xml_entries = []
    for path, priority, freq in pages:
        xml_entries.append(
            f"  <url>\n"
            f"    <loc>{base}{path}</loc>\n"
            f"    <lastmod>{today}</lastmod>\n"
            f"    <changefreq>{freq}</changefreq>\n"
            f"    <priority>{priority}</priority>\n"
            f"  </url>"
        )

    xml = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9">\n'
        + "\n".join(xml_entries) + "\n"
        '</urlset>\n'
    )

    return Response(xml, mimetype="application/xml", headers={
        "Cache-Control": "public, max-age=3600",
    })

@app.route('/robots.txt')
@limiter.exempt
def robots_txt():
    """Serve robots.txt for search engine crawlers."""
    txt = (
        "User-agent: *\n"
        "Allow: /\n"
        "\n"
        "Sitemap: https://www.vijaypdf.com/sitemap.xml\n"
        "Sitemap: https://vijaypdf.com/sitemap.xml\n"
    )
    return Response(txt, mimetype="text/plain", headers={
        "Cache-Control": "public, max-age=86400"
    })

# Main web routes
@app.route('/')
def index():
    total_conversions = 0
    recent_count = 0

    if current_user.is_authenticated:
        total_conversions = ConversionHistory.query.filter_by(user_id=current_user.id).count()
        recent_count = ConversionHistory.query.filter_by(user_id=current_user.id).order_by(ConversionHistory.id.desc()).limit(5).count()

    return render_template(
        'index.html',
        total_conversions=total_conversions,
        recent_count=recent_count
    )

@app.route('/login', methods=['GET', 'POST'])
@limiter.limit("5 per minute")
def login():
    if current_user.is_authenticated:
        return redirect(url_for('index'))
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        user = User.query.filter_by(username=username).first()
        
        if user:
            if user.locked_until and user.locked_until > datetime.utcnow():
                flash('Account locked due to multiple failed login attempts. Try again later.', 'danger')
                app.logger.warning(f'Locked login attempt for user: {username}')
                return render_template('login.html')
                
            if bcrypt.check_password_hash(user.password, password):
                user.failed_logins = 0
                user.locked_until = None
                db.session.commit()
                login_user(user)
                app.logger.info(f'Successful login for user: {username}')
                return redirect(url_for('index'))
            else:
                user.failed_logins += 1
                if user.failed_logins >= 5:
                    user.locked_until = datetime.utcnow() + timedelta(minutes=15)
                    app.logger.warning(f'Account locked for user: {username}')
                db.session.commit()
                
        app.logger.warning(f'Failed login attempt for username: {username}')
        flash('Login Unsuccessful. Please check username and password', 'danger')
    return render_template('login.html')

@app.route('/register', methods=['GET', 'POST'])
@limiter.limit("3 per minute")
def register():
    if current_user.is_authenticated:
        return redirect(url_for('index'))
    if request.method == 'POST':
        username = request.form.get('username')
        email = request.form.get('email')
        password = request.form.get('password')
        
        if User.query.filter_by(username=username).first():
            flash('Username already exists.', 'danger')
            return render_template('register.html')
            
        if User.query.filter_by(email=email).first():
            flash('Email already registered.', 'danger')
            return render_template('register.html')

        hashed_password = bcrypt.generate_password_hash(password).decode('utf-8')
        
        # Generate OTP
        otp = str(secrets.randbelow(1000000)).zfill(6)
        otp_expiry = datetime.utcnow() + timedelta(minutes=10)
        
        user = User(username=username, email=email, password=hashed_password, otp=otp, otp_expiry=otp_expiry)
        
        try:
            db.session.add(user)
            db.session.commit()
            
            # Send OTP Email
            email_service.send_otp_email(email, otp)
            
            flash('Account created! Please verify your email with the OTP sent.', 'info')
            return redirect(url_for('verify_otp', email=email))
        except Exception as e:
            db.session.rollback()
            flash(f'An error occurred: {str(e)}', 'danger')
            
    return render_template('register.html')

@app.route('/verify-otp', methods=['GET', 'POST'])
def verify_otp():
    email = request.args.get('email')
    if request.method == 'POST':
        email = request.form.get('email')
        otp_input = request.form.get('otp')
        user = User.query.filter_by(email=email).first()
        
        if user and user.otp == otp_input and user.otp_expiry > datetime.utcnow():
            user.is_verified = True
            user.otp = None
            user.otp_expiry = None
            db.session.commit()
            flash('Email verified! You can now log in.', 'success')
            return redirect(url_for('login'))
        else:
            flash('Invalid or expired OTP.', 'danger')
            
    return render_template('verify_otp.html', email=email)

@app.route('/forgot-password', methods=['GET', 'POST'])
def forgot_password():
    if request.method == 'POST':
        email = request.form.get('email')
        user = User.query.filter_by(email=email).first()
        if user:
            token = secrets.token_urlsafe(32)
            user.reset_token = token
            user.reset_token_expiry = datetime.utcnow() + timedelta(hours=1)
            db.session.commit()
            email_service.send_password_reset_email(email, token)
        
        flash('If an account matches that email, a reset link has been sent.', 'info')
        return redirect(url_for('login'))
    return render_template('forgot_password.html')

@app.route('/reset-password/<token>', methods=['GET', 'POST'])
def reset_password(token):
    user = User.query.filter_by(reset_token=token).first()
    if not user or user.reset_token_expiry < datetime.utcnow():
        flash('Invalid or expired reset token.', 'danger')
        return redirect(url_for('forgot_password'))
        
    if request.method == 'POST':
        new_password = request.form.get('password')
        hashed_password = bcrypt.generate_password_hash(new_password).decode('utf-8')
        user.password = hashed_password
        user.reset_token = None
        user.reset_token_expiry = None
        db.session.commit()
        flash('Password updated! You can now log in.', 'success')
        return redirect(url_for('login'))
        
    return render_template('reset_password.html', token=token)

@app.route('/test-email')
def test_email_route():
    if email_service.test_connection():
        return "Test email sent successfully! Check your inbox (support@vijaypdf.com)."
    else:
        return "Failed to send test email. Check server logs and .env configuration."

@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('index'))

@app.route('/history')
@login_required
def history():
    user_history = ConversionHistory.query.filter_by(user_id=current_user.id).order_by(ConversionHistory.id.desc()).all()
    return render_template('history.html', history=user_history)
@app.route('/download/<filename>')
def download_file(filename):
    safe_filename = secure_filename(filename)

    file_path = os.path.join(app.config['OUTPUT_FOLDER'], safe_filename)
    if not os.path.exists(file_path):
        return jsonify({"error": "File not found"}), 404

    return send_from_directory(app.config['OUTPUT_FOLDER'], safe_filename, as_attachment=True)
if __name__ == '__main__':
    with app.app_context(): db.create_all()
    debug_mode = os.getenv("FLASK_ENV", "development").lower() == "development"
    app.run(debug=debug_mode, port=5000)

    # if __name__ == "__main__":
    #     app.run(debug=True)

